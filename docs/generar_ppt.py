# -*- coding: utf-8 -*-
"""
Genera la presentacion (PPTX) del proyecto Compilador ProyLang.
Autor: Cesar Daniel Osorio Polanco
Ejecutar:  python docs\generar_ppt.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- Paleta de colores ----
AZUL    = RGBColor(0x1F, 0x3A, 0x5F)   # azul oscuro (titulos)
AZUL2   = RGBColor(0x2E, 0x6D, 0xB4)   # azul medio
GRIS    = RGBColor(0x33, 0x33, 0x33)   # texto
VERDE   = RGBColor(0x2E, 0x8B, 0x57)
BLANCO  = RGBColor(0xFF, 0xFF, 0xFF)
GRISCLR = RGBColor(0xF0, 0xF3, 0xF7)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def add_slide(titulo):
    s = prs.slides.add_slide(BLANK)
    # banda de titulo
    banda = s.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.1))
    banda.fill.solid(); banda.fill.fore_color.rgb = AZUL
    banda.line.fill.background()
    tf = banda.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.4); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = titulo
    p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = BLANCO
    return s

def add_bullets(s, items, top=1.4, left=0.7, width=12.0, size=20):
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.6))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (txt, lvl) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt; p.level = lvl
        p.font.size = Pt(size - lvl*2); p.font.color.rgb = GRIS
        p.space_after = Pt(8)
        if lvl == 0:
            p.font.bold = True; p.font.color.rgb = AZUL2
    return s

def add_code(s, codigo, top=1.5, left=0.8, width=11.7, height=4.8, size=15):
    box = s.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid(); box.fill.fore_color.rgb = GRISCLR
    box.line.color.rgb = AZUL2; box.line.width = Pt(1)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.2)
    for i, linea in enumerate(codigo.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = linea; p.font.name = "Consolas"; p.font.size = Pt(size)
        p.font.color.rgb = GRIS
    return s

# ============================================================
# Slide 1 - Portada
# ============================================================
s = prs.slides.add_slide(BLANK)
fondo = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
fondo.fill.solid(); fondo.fill.fore_color.rgb = AZUL; fondo.line.fill.background()
tb = s.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.3), Inches(3))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "Compilador para un DSL de"
p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = BLANCO; p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph(); p.text = "Gestion de Proyectos  (ProyLang)"
p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = BLANCO; p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph(); p.text = ""
p = tf.add_paragraph(); p.text = "Cesar Daniel Osorio Polanco   |   Carne 0900-16-6607"
p.font.size = Pt(20); p.font.color.rgb = RGBColor(0xCF, 0xDC, 0xEC); p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph(); p.text = "ANTLR 4  +  Java 17"
p.font.size = Pt(18); p.font.color.rgb = RGBColor(0xCF, 0xDC, 0xEC); p.alignment = PP_ALIGN.CENTER

# ============================================================
# Slide 2 - Objetivo
# ============================================================
s = add_slide("Objetivo del proyecto")
add_bullets(s, [
    ("Disenar e implementar un compilador funcional para un lenguaje propio (DSL).", 0),
    ("Aplicar las cuatro fases clasicas de un compilador:", 0),
    ("Analisis lexico  -  convertir el texto en tokens", 1),
    ("Analisis sintactico  -  verificar la gramatica", 1),
    ("Analisis semantico  -  validar el significado (tabla de simbolos)", 1),
    ("Generacion de codigo  -  producir la salida", 1),
    ("Detectar y reportar errores en cada fase.", 0),
])

# ============================================================
# Slide 3 - Que es ProyLang
# ============================================================
s = add_slide("El lenguaje: ProyLang")
add_bullets(s, [
    ("Es un DSL (lenguaje de dominio especifico) para describir proyectos.", 0),
    ("Permite definir tareas con sus atributos:", 0),
    ("duracion (obligatoria), costo, responsable, prioridad", 1),
    ("dependencias entre tareas (que va antes de que)", 1),
    ("A partir de esa descripcion, el compilador calcula automaticamente:", 0),
    ("el cronograma del proyecto y su duracion total", 1),
    ("la RUTA CRITICA (metodo CPM) y la holgura de cada tarea", 1),
])

# ============================================================
# Slide 4 - Ejemplo de codigo
# ============================================================
s = add_slide("Ejemplo de un programa en ProyLang")
add_code(s, '''proyecto "App Movil" {

    tarea diseno {
        duracion    = 5;
        costo       = 5 * 800;
        responsable = "Ana Lopez";
        prioridad   = alta;
    }

    tarea backend {
        duracion = 10;
        depende  = diseno;
    }
}''')

# ============================================================
# Slide 5 - Arquitectura (4 fases)
# ============================================================
s = add_slide("Arquitectura: las 4 fases")
add_code(s, '''   archivo.proy
        |
        v
  [ LEXER ]  --tokens-->  [ PARSER ]  --arbol-->  [ SEMANTICO ]
   Fase 1                  Fase 2                   Fase 3
                                                      |
                                                      v
                                          [ GENERADOR DE CODIGO ]
                                                  Fase 4
                                    +-------------+-------------+
                                    v             v             v
                              3 direcciones    JSON (CPM)     Plan

  Si una fase encuentra errores, el proceso se detiene.''', size=14, height=5.0)

# ============================================================
# Slide 6 - Fase 1 lexico
# ============================================================
s = add_slide("Fase 1 - Analisis Lexico")
add_bullets(s, [
    ("Convierte el texto en TOKENS (las 'palabras' del lenguaje).", 0),
    ("Ejemplo:  duracion = 5;  produce:", 0),
    ("KW_DURACION  ->  'duracion'", 1),
    ("ASIGNA       ->  '='", 1),
    ("NUMERO       ->  '5'", 1),
    ("PYC          ->  ';'", 1),
    ("Se ignoran espacios y comentarios.", 0),
    ("Error lexico: un caracter que no pertenece al lenguaje (ej. '@', '#').", 0),
])

# ============================================================
# Slide 7 - Fase 2 sintactico
# ============================================================
s = add_slide("Fase 2 - Analisis Sintactico")
add_bullets(s, [
    ("Verifica que los tokens cumplan la GRAMATICA del lenguaje.", 0),
    ("La gramatica esta definida en el archivo ProyLang.g4 (ANTLR).", 0),
    ("Construye el arbol de sintaxis (estructura del programa).", 0),
    ("Comprueba: llaves { }, punto y coma ;, estructura de tareas, etc.", 0),
    ("Error sintactico: falta un ';', una llave '}', orden incorrecto...", 0),
])

# ============================================================
# Slide 8 - Fase 3 semantico
# ============================================================
s = add_slide("Fase 3 - Analisis Semantico")
add_bullets(s, [
    ("Construye la TABLA DE SIMBOLOS (cada tarea con sus datos).", 0),
    ("Verifica que el programa tenga SENTIDO. Validaciones:", 0),
    ("Tareas duplicadas / atributos repetidos", 1),
    ("'duracion' obligatoria y mayor que 0; 'costo' no negativo", 1),
    ("'prioridad' debe ser alta / media / baja", 1),
    ("Dependencias deben apuntar a tareas existentes (no a si misma)", 1),
    ("Deteccion de DEPENDENCIAS CIRCULARES (DFS con coloreo)", 1),
])

# ============================================================
# Slide 9 - Fase 4 generacion
# ============================================================
s = add_slide("Fase 4 - Generacion de Codigo")
add_bullets(s, [
    ("Traduccion dirigida por la sintaxis. Produce 3 salidas:", 0),
    ("Codigo intermedio de TRES DIRECCIONES (de las expresiones)", 1),
    ("ej: costo = 10*1000+500  ->  t1 = 10*1000 ; t2 = t1+500", 1),
    ("Cronograma en JSON con el metodo de la RUTA CRITICA (CPM)", 1),
    ("inicio/fin temprano y tardio, holgura, duracion total", 1),
    ("Plan legible con la ejecucion simulada dia por dia", 1),
])

# ============================================================
# Slide 10 - Manejo de errores
# ============================================================
s = add_slide("Manejo de errores en cada fase")
add_bullets(s, [
    ("El compilador reporta los errores con su FASE, LINEA y COLUMNA.", 0),
    ("Tres tipos de error, cada uno con su ejemplo de prueba:", 0),
    ("Lexico   -> error_lexico.proy   (caracteres invalidos @ #)", 1),
    ("Sintactico -> error_sintactico.proy  (falta ; y })", 1),
    ("Semantico -> error_semantico.proy  (8 errores distintos)", 1),
    ("Si una fase falla, NO se pasa a la siguiente (como un compilador real).", 0),
])

# ============================================================
# Slide 11 - Resultado / demo
# ============================================================
s = add_slide("Resultado: cronograma y ruta critica")
add_code(s, ''' PLAN DEL PROYECTO: App Movil
 Tareas: 5   |   Duracion total: 21 dias

 TAREA        DUR   INI   FIN   HOLG  CRIT
 -----------------------------------------
 diseno         5     0     5     0   SI
 backend       10     5    15     0   SI
 frontend       8     5    13     2
 pruebas        4    15    19     0   SI
 despliegue     2    19    21     0   SI

 Ruta critica: diseno -> backend -> pruebas -> despliegue''', size=15, height=4.6)

# ============================================================
# Slide 12 - Tecnologias
# ============================================================
s = add_slide("Tecnologias utilizadas")
add_bullets(s, [
    ("ANTLR 4.13.2  -  generacion del lexer y el parser desde la gramatica .g4", 0),
    ("Java 17  -  analisis semantico y generacion de codigo", 0),
    ("Git + GitHub  -  control de versiones del proyecto", 0),
    ("Repositorio:", 0),
    ("github.com/CesarDanielOsorio/compilador-dsl-proylang", 1),
])

# ============================================================
# Slide 13 - Conclusiones
# ============================================================
s = add_slide("Conclusiones")
add_bullets(s, [
    ("Se implemento un compilador funcional completo con las 4 fases.", 0),
    ("El DSL resuelve un problema real: planificar proyectos y su ruta critica.", 0),
    ("Cada fase maneja sus propios errores de forma clara.", 0),
    ("El diseno es escalable: agregar atributos sigue siempre el mismo patron.", 0),
])

# ============================================================
# Slide 14 - Gracias
# ============================================================
s = prs.slides.add_slide(BLANK)
fondo = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
fondo.fill.solid(); fondo.fill.fore_color.rgb = AZUL; fondo.line.fill.background()
tb = s.shapes.add_textbox(Inches(1), Inches(3), Inches(11.3), Inches(2))
tf = tb.text_frame
p = tf.paragraphs[0]; p.text = "Gracias"
p.font.size = Pt(54); p.font.bold = True; p.font.color.rgb = BLANCO; p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph(); p.text = "Cesar Daniel Osorio Polanco"
p.font.size = Pt(22); p.font.color.rgb = RGBColor(0xCF, 0xDC, 0xEC); p.alignment = PP_ALIGN.CENTER

prs.save("docs/Presentacion_ProyLang.pptx")
print("PPTX generado: docs/Presentacion_ProyLang.pptx  (" + str(len(prs.slides.__iter__.__self__._sldIdLst)) + " slides)")
